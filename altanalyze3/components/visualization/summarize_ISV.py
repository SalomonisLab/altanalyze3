#!/usr/bin/env python3
"""Build a PowerPoint summary from ISV output folders.

Expected input structure:
- <results_dir>/<GENE>/
    - combined__...pdf
    - <condition>__..._isoform_ids.tsv

For each gene directory, this script:
1) Finds a combined PDF (prefers *_custom.pdf when available).
2) Converts the first page to an image (JPG or PNG).
3) Creates one PPT slide:
   - gene name centered near top (font 20)
   - image in the center area
   - bottom text with top isoforms per cluster ID from TSVs.
"""

from __future__ import annotations

import argparse
import csv
import glob
import math
import os
import shutil
import subprocess
from collections import defaultdict
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple


def _import_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        return Presentation, Inches, Pt
    except Exception as exc:
        raise RuntimeError(
            "python-pptx is required. Install with: pip install python-pptx"
        ) from exc


@dataclass
class GeneInput:
    gene: str
    gene_dir: str
    combined_pdf: str


@dataclass
class GeneScore:
    score: float
    min_condition_reads: float
    major_clusters: Tuple[str, str]
    cluster_shifts: Tuple[float, float]
    raw_ratio_shift: float
    read_weight: float
    weight_factor: float


def find_gene_inputs(results_dir: str, combined_glob: str = "combined__*.pdf") -> List[GeneInput]:
    gene_inputs: List[GeneInput] = []
    for name in sorted(os.listdir(results_dir)):
        gene_dir = os.path.join(results_dir, name)
        if not os.path.isdir(gene_dir):
            continue
        matches = glob.glob(os.path.join(gene_dir, combined_glob))
        if not matches:
            continue
        # Prefer custom render when available.
        custom = sorted([m for m in matches if m.endswith("_custom.pdf")])
        selected = custom[0] if custom else sorted(matches)[0]
        gene_inputs.append(GeneInput(gene=name, gene_dir=gene_dir, combined_pdf=selected))
    return gene_inputs


def _convert_pdf_with_pymupdf(pdf_path: str, image_path: str, dpi: int) -> bool:
    try:
        import fitz  # type: ignore
    except Exception:
        return False
    zoom = max(dpi, 72) / 72.0
    doc = fitz.open(pdf_path)
    try:
        if len(doc) == 0:
            return False
        page = doc.load_page(0)
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        pix.save(image_path)
        return True
    finally:
        doc.close()


def _convert_pdf_with_pdf2image(pdf_path: str, image_path: str, dpi: int, image_format: str) -> bool:
    try:
        from pdf2image import convert_from_path  # type: ignore
    except Exception:
        return False
    try:
        images = convert_from_path(pdf_path, dpi=dpi, first_page=1, last_page=1)
        if not images:
            return False
        fmt = "PNG" if image_format == "png" else "JPEG"
        save_kwargs = {"format": fmt}
        if image_format == "jpg":
            save_kwargs["quality"] = 95
        images[0].convert("RGB").save(image_path, **save_kwargs)
        return True
    except Exception:
        # Common when poppler/pdfinfo is missing from PATH.
        return False


def _convert_pdf_with_pypdfium2(pdf_path: str, image_path: str, dpi: int, image_format: str) -> bool:
    try:
        import pypdfium2 as pdfium  # type: ignore
    except Exception:
        return False
    try:
        doc = pdfium.PdfDocument(pdf_path)
        if len(doc) == 0:
            return False
        page = doc[0]
        scale = max(dpi, 72) / 72.0
        bitmap = page.render(scale=scale)
        image = bitmap.to_pil()
        fmt = "PNG" if image_format == "png" else "JPEG"
        save_kwargs = {"format": fmt}
        if image_format == "jpg":
            save_kwargs["quality"] = 95
        image.convert("RGB").save(image_path, **save_kwargs)
        return True
    except Exception:
        return False


def _convert_pdf_with_pdftoppm(pdf_path: str, image_path: str, dpi: int, image_format: str) -> bool:
    tool = shutil.which("pdftoppm")
    if not tool:
        return False
    out_prefix = os.path.splitext(image_path)[0]
    ppm_flag = "-png" if image_format == "png" else "-jpeg"
    cmd = [
        tool,
        ppm_flag,
        "-f",
        "1",
        "-singlefile",
        "-r",
        str(int(dpi)),
        pdf_path,
        out_prefix,
    ]
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except Exception:
        return False
    return os.path.exists(image_path)


def convert_pdf_to_image(
    pdf_path: str,
    image_path: str,
    dpi: int = 200,
    image_format: str = "jpg",
) -> None:
    image_format = str(image_format).lower()
    if image_format not in {"jpg", "png"}:
        raise ValueError("image_format must be 'jpg' or 'png'")
    os.makedirs(os.path.dirname(image_path), exist_ok=True)
    if _convert_pdf_with_pymupdf(pdf_path, image_path, dpi):
        return
    if _convert_pdf_with_pypdfium2(pdf_path, image_path, dpi, image_format):
        return
    if _convert_pdf_with_pdf2image(pdf_path, image_path, dpi, image_format):
        return
    if _convert_pdf_with_pdftoppm(pdf_path, image_path, dpi, image_format):
        return
    raise RuntimeError(
        f"Could not convert PDF: {pdf_path}\n"
        "Install one of: pymupdf (fitz), pypdfium2, pdf2image (+ poppler), or pdftoppm.\n"
        "For pip-only installs: pip3 install pypdfium2\n"
        "If using pdf2image on macOS: brew install poppler"
    )


def convert_pdf_to_jpg(pdf_path: str, jpg_path: str, dpi: int = 200) -> None:
    # Backward-compatible wrapper.
    convert_pdf_to_image(pdf_path, jpg_path, dpi=dpi, image_format="jpg")


def _safe_float(value: str, default: float = 0.0) -> float:
    try:
        return float(value)
    except Exception:
        return default


def _structure_length_metric(structure: str) -> int:
    if not structure:
        return 0
    return sum(1 for tok in structure.split("|") if tok.strip())


def _condition_from_tsv_path(path: str) -> Optional[str]:
    base = os.path.basename(path)
    if base.startswith("combined__"):
        return None
    if "__" not in base:
        return None
    return base.split("__", 1)[0]


def _condition_read_totals(gene_dir: str) -> Dict[str, float]:
    totals: Dict[str, float] = defaultdict(float)
    tsv_paths = sorted(glob.glob(os.path.join(gene_dir, "*_isoform_ids.tsv")))
    for path in tsv_paths:
        condition = _condition_from_tsv_path(path)
        if not condition:
            continue
        with open(path, "r", newline="") as handle:
            reader = csv.DictReader(handle, delimiter="\t")
            if not reader.fieldnames:
                continue
            for row in reader:
                count = _safe_float(row.get("count", "0"), 0.0)
                if count <= 0:
                    count = 1.0
                totals[condition] += count
    return totals


def _passes_min_reads_any_condition(gene_dir: str, min_reads: int) -> bool:
    if min_reads <= 0:
        return True
    totals = _condition_read_totals(gene_dir)
    if not totals:
        return False
    return any(total >= float(min_reads) for total in totals.values())


def _safe_int(value: str, default: int = 0) -> int:
    try:
        return int(float(value))
    except Exception:
        return default


def _normalize_transcript_id(label: str) -> str:
    tx = str(label or "").strip()
    if "|" in tx:
        tx = tx.split("|", 1)[0].strip()
    tx = tx.replace("(NMD)", "").strip()
    return tx


def _tx_id_candidates(label: str) -> List[str]:
    tx = _normalize_transcript_id(label)
    if not tx:
        return []
    candidates = [tx]
    if tx.startswith("ENST") and "." in tx:
        candidates.append(tx.split(".", 1)[0])
    return candidates


def _build_isoform_cluster_lookup(gene_dir: str) -> Dict[str, str]:
    tsv_paths = sorted(glob.glob(os.path.join(gene_dir, "*_isoform_ids.tsv")))
    tx_cluster_counts: Dict[str, Dict[str, float]] = defaultdict(lambda: defaultdict(float))

    for path in tsv_paths:
        with open(path, "r", newline="") as handle:
            reader = csv.DictReader(handle, delimiter="\t")
            if not reader.fieldnames:
                continue
            for row in reader:
                cluster = (row.get("cluster_label") or row.get("cluster_index") or "").strip()
                if not cluster:
                    continue
                count = _safe_float(row.get("count", "0"), 0.0)
                if count <= 0:
                    count = 1.0
                tx_candidates = [
                    (row.get("resolved_id") or "").strip(),
                    (row.get("isoform_id") or "").strip(),
                ]
                for tx in tx_candidates:
                    if not tx:
                        continue
                    norm = _normalize_transcript_id(tx)
                    if not norm:
                        continue
                    tx_cluster_counts[norm][cluster] += count
                    if norm.startswith("ENST") and "." in norm:
                        tx_cluster_counts[norm.split(".", 1)[0]][cluster] += count

    lookup: Dict[str, str] = {}
    for tx, cluster_counts in tx_cluster_counts.items():
        if not cluster_counts:
            continue
        lookup[tx] = max(cluster_counts.items(), key=lambda kv: kv[1])[0]
    return lookup


def _load_event_rows(event_table_tsv: str) -> Tuple[Dict[str, List[Dict[str, Any]]], List[str]]:
    by_gene: Dict[str, List[Dict[str, Any]]] = defaultdict(list)
    diff_cols: List[str] = []
    with open(event_table_tsv, "r", newline="") as handle:
        reader = csv.DictReader(handle, delimiter="\t")
        if not reader.fieldnames:
            return by_gene, diff_cols
        diff_cols = [c for c in reader.fieldnames if c.startswith("DiffMeans_")]
        for row in reader:
            gene = (row.get("Gene_Symbol") or "").strip()
            feature = (row.get("Feature") or "").strip()
            if not gene or not feature:
                continue
            diff_map: Dict[str, str] = {}
            for col in diff_cols:
                value = (row.get(col) or "").strip()
                if value:
                    diff_map[col] = value
            if diff_map:
                max_abs = max(abs(_safe_float(v, 0.0)) for v in diff_map.values())
            else:
                max_abs = 0.0
            by_gene[gene].append(
                {
                    "feature": feature,
                    "coord_1": (row.get("Coord_1") or "").strip(),
                    "coord_2": (row.get("Coord_2") or "").strip(),
                    "isoform_1": (row.get("Isoform_1|Length") or "").strip(),
                    "isoform_2": (row.get("Isoform_2|Length") or "").strip(),
                    "diff_map": diff_map,
                    "event_type": (row.get("Event_Type") or "").strip(),
                    "event_direction": (row.get("Event_Direction") or "").strip(),
                    "max_abs_diff": max_abs,
                }
            )
    for gene, rows in by_gene.items():
        rows.sort(key=lambda r: (r.get("max_abs_diff", 0.0), r.get("feature", "")), reverse=True)
    return by_gene, diff_cols


def _render_event_summary(
    rows: List[Dict[str, Any]],
    isoform_cluster_lookup: Dict[str, str],
    max_events_per_gene: int = 3,
) -> str:
    if not rows:
        return ""
    limit = max_events_per_gene if max_events_per_gene and max_events_per_gene > 0 else len(rows)
    lines: List[str] = ["Matched events from provided table:"]
    for row in rows[:limit]:
        iso1 = row.get("isoform_1", "")
        iso2 = row.get("isoform_2", "")
        iso1_cluster = "NA"
        for candidate in _tx_id_candidates(iso1):
            if candidate in isoform_cluster_lookup:
                iso1_cluster = isoform_cluster_lookup[candidate]
                break
        iso2_cluster = "NA"
        for candidate in _tx_id_candidates(iso2):
            if candidate in isoform_cluster_lookup:
                iso2_cluster = isoform_cluster_lookup[candidate]
                break
        diff_map = row.get("diff_map", {}) or {}
        if diff_map:
            diff_txt = "; ".join(f"{k}={v}" for k, v in diff_map.items())
        else:
            diff_txt = "DiffMeans=NA"
        coord1 = row.get("coord_1", "") or "NA"
        coord2 = row.get("coord_2", "") or "NA"
        lines.append(
            f"{row.get('feature','')} | coords={coord1},{coord2} | "
            f"iso1={iso1}[cluster={iso1_cluster}] | iso2={iso2}[cluster={iso2_cluster}] | {diff_txt}"
        )
    return "\n".join(lines)


def _exon_tokens_in_order(structure: str) -> List[str]:
    seen = set()
    ordered: List[str] = []
    for tok in structure.split("|"):
        token = tok.strip()
        if not token or not token.startswith("E"):
            continue
        if token in seen:
            continue
        seen.add(token)
        ordered.append(token)
    return ordered


def _short_unique_exon_label(gene_dir: str, max_each: int = 2) -> str:
    """Return concise label of exons unique between top 2 clusters.

    Format examples:
    - "E1.2_1234 vs E1.1_1200"
    - "E1.2_1234,E2.1 (+1) vs E1.1_1200"
    """
    tsv_paths = sorted(glob.glob(os.path.join(gene_dir, "*_isoform_ids.tsv")))
    cluster_totals: Dict[str, float] = defaultdict(float)
    cluster_structure_counts: Dict[str, Dict[str, float]] = defaultdict(lambda: defaultdict(float))

    for path in tsv_paths:
        with open(path, "r", newline="") as handle:
            reader = csv.DictReader(handle, delimiter="\t")
            if not reader.fieldnames:
                continue
            for row in reader:
                cluster = (row.get("cluster_label") or row.get("cluster_index") or "").strip()
                if not cluster:
                    continue
                structure = (
                    row.get("tokens_trimmed")
                    or row.get("cluster_tokens_exon_intron")
                    or row.get("cluster_tokens_exon")
                    or row.get("tokens_raw")
                    or ""
                ).strip()
                if not structure:
                    continue
                count = _safe_float(row.get("count", "0"), 0.0)
                if count <= 0:
                    count = 1.0
                cluster_totals[cluster] += count
                cluster_structure_counts[cluster][structure] += count

    if len(cluster_totals) < 2:
        return ""

    top_clusters = sorted(cluster_totals.items(), key=lambda kv: kv[1], reverse=True)[:2]
    cluster_a, cluster_b = top_clusters[0][0], top_clusters[1][0]
    struct_a = max(cluster_structure_counts[cluster_a].items(), key=lambda kv: kv[1])[0]
    struct_b = max(cluster_structure_counts[cluster_b].items(), key=lambda kv: kv[1])[0]

    exons_a = _exon_tokens_in_order(struct_a)
    exons_b = _exon_tokens_in_order(struct_b)
    set_b = set(exons_b)
    set_a = set(exons_a)
    uniq_a = [x for x in exons_a if x not in set_b]
    uniq_b = [x for x in exons_b if x not in set_a]

    def _fmt(tokens: List[str]) -> str:
        if not tokens:
            return "none"
        shown = tokens[:max_each]
        extra = len(tokens) - len(shown)
        txt = ",".join(shown)
        if extra > 0:
            txt += f" (+{extra})"
        return txt

    return f"{_fmt(uniq_a)} vs {_fmt(uniq_b)}"


def compute_gene_shift_score(gene_dir: str) -> GeneScore:
    """Score a gene by the largest two-cluster ratio shift across conditions.

    Steps:
    1) Find the 2 major clusters by total weighted reads across all conditions.
    2) For each major cluster, compute ratio per condition:
       ratio = cluster_reads_in_condition / total_reads_in_condition
    3) Per cluster, compute ratio shift (log2 fold-change):
       log2((max_ratio + eps) / (min_ratio + eps))
    4) Weight by minimum condition depth with a bounded factor:
       w = n / (n + K), factor = base + span*w
    5) Gene score = max(cluster_shift_1, cluster_shift_2) * factor
    """
    tsv_paths = sorted(glob.glob(os.path.join(gene_dir, "*_isoform_ids.tsv")))
    condition_cluster_counts: Dict[str, Dict[str, float]] = defaultdict(lambda: defaultdict(float))
    condition_totals: Dict[str, float] = defaultdict(float)
    global_cluster_totals: Dict[str, float] = defaultdict(float)

    for path in tsv_paths:
        condition = _condition_from_tsv_path(path)
        if not condition:
            continue
        with open(path, "r", newline="") as handle:
            reader = csv.DictReader(handle, delimiter="\t")
            if not reader.fieldnames:
                continue
            for row in reader:
                cluster = (row.get("cluster_label") or row.get("cluster_index") or "").strip()
                if not cluster:
                    continue
                count = _safe_float(row.get("count", "0"), 0.0)
                if count <= 0:
                    count = 1.0
                condition_cluster_counts[condition][cluster] += count
                condition_totals[condition] += count
                global_cluster_totals[cluster] += count

    if len(global_cluster_totals) < 2 or len(condition_totals) < 2:
        return GeneScore(
            score=0.0,
            min_condition_reads=0.0,
            major_clusters=("NA", "NA"),
            cluster_shifts=(0.0, 0.0),
            raw_ratio_shift=0.0,
            read_weight=0.0,
            weight_factor=0.0,
        )

    top_clusters = sorted(
        global_cluster_totals.items(),
        key=lambda kv: kv[1],
        reverse=True,
    )[:2]
    cluster_a, cluster_b = top_clusters[0][0], top_clusters[1][0]

    def _cluster_shift(cluster_id: str) -> float:
        ratios: List[float] = []
        for condition, total_reads in condition_totals.items():
            if total_reads <= 0:
                continue
            ratios.append(condition_cluster_counts[condition].get(cluster_id, 0.0) / total_reads)
        if not ratios:
            return 0.0
        max_ratio = max(ratios)
        min_ratio = min(ratios)
        eps = 0.01
        return math.log2((max_ratio + eps) / (min_ratio + eps))

    shift_a = _cluster_shift(cluster_a)
    shift_b = _cluster_shift(cluster_b)
    min_condition_reads = min(v for v in condition_totals.values() if v > 0) if condition_totals else 0.0
    raw_ratio_shift = max(shift_a, shift_b)
    k_reads = 100.0
    weight_base = 0.75
    weight_span = 0.25
    read_weight = (min_condition_reads / (min_condition_reads + k_reads)) if min_condition_reads > 0 else 0.0
    weight_factor = weight_base + (weight_span * read_weight)
    score = raw_ratio_shift * weight_factor

    return GeneScore(
        score=score,
        min_condition_reads=min_condition_reads,
        major_clusters=(cluster_a, cluster_b),
        cluster_shifts=(shift_a, shift_b),
        raw_ratio_shift=raw_ratio_shift,
        read_weight=read_weight,
        weight_factor=weight_factor,
    )


def _collect_cluster_summaries(gene_dir: str) -> List[Dict[str, Any]]:
    tsv_paths = sorted(glob.glob(os.path.join(gene_dir, "*_isoform_ids.tsv")))
    # In practice combined files may or may not have an associated tsv; include all TSVs found.
    if not tsv_paths:
        return []

    cluster_totals: Dict[str, float] = defaultdict(float)
    cluster_structure_counts: Dict[str, Dict[str, float]] = defaultdict(lambda: defaultdict(float))
    structure_tx_counts: Dict[Tuple[str, str], Dict[str, float]] = defaultdict(lambda: defaultdict(float))
    cluster_tx_lengths: Dict[str, Dict[str, int]] = defaultdict(dict)

    for path in tsv_paths:
        with open(path, "r", newline="") as handle:
            reader = csv.DictReader(handle, delimiter="\t")
            if not reader.fieldnames:
                continue
            for row in reader:
                cluster = (row.get("cluster_label") or row.get("cluster_index") or "").strip()
                if not cluster:
                    continue
                structure = (
                    row.get("tokens_trimmed")
                    or row.get("cluster_tokens_exon_intron")
                    or row.get("cluster_tokens_exon")
                    or row.get("tokens_raw")
                    or ""
                ).strip()
                if not structure:
                    continue
                tx_id = (row.get("resolved_id") or row.get("isoform_id") or "").strip()
                if not tx_id:
                    tx_id = "NA"
                count = _safe_float(row.get("count", "0"), 0.0)
                if count <= 0:
                    count = 1.0
                length_metric = _structure_length_metric(structure)
                cluster_totals[cluster] += count
                cluster_structure_counts[cluster][structure] += count
                structure_tx_counts[(cluster, structure)][tx_id] += count
                prev_len = cluster_tx_lengths[cluster].get(tx_id, 0)
                if length_metric > prev_len:
                    cluster_tx_lengths[cluster][tx_id] = length_metric

    if not cluster_totals:
        return []

    ordered_clusters = sorted(cluster_totals.items(), key=lambda kv: kv[1], reverse=True)

    rows: List[Dict[str, Any]] = []
    for cluster, _total in ordered_clusters:
        structures = cluster_structure_counts.get(cluster, {})
        if not structures:
            continue
        structure, structure_count = max(structures.items(), key=lambda kv: kv[1])
        tx_counts = structure_tx_counts.get((cluster, structure), {})
        top_tx, top_tx_count = ("NA", 0.0)
        if tx_counts:
            # Keep representative transcript frequency-driven, but restrict
            # candidates to the top 10% longest reads in the cluster.
            tx_len_map = cluster_tx_lengths.get(cluster, {})
            cluster_lengths = sorted(tx_len_map.values(), reverse=True)
            eligible_tx_counts = tx_counts
            if cluster_lengths:
                keep_n = max(1, int(math.ceil(len(cluster_lengths) * 0.10)))
                min_top_length = cluster_lengths[keep_n - 1]
                constrained = {
                    tx: c for tx, c in tx_counts.items()
                    if tx_len_map.get(tx, 0) >= min_top_length
                }
                if constrained:
                    eligible_tx_counts = constrained
            top_tx, top_tx_count = max(
                eligible_tx_counts.items(),
                key=lambda kv: (kv[1], tx_len_map.get(kv[0], 0), kv[0]),
            )
        rows.append(
            {
                "cluster": cluster,
                "structure": structure,
                "transcript_id": top_tx,
                "structure_n": int(structure_count),
                "tx_n": int(top_tx_count),
                "cluster_total": int(cluster_totals.get(cluster, 0)),
            }
        )
    return rows


def summarize_clusters_from_tsvs(
    gene_dir: str,
    top_n: int = 0,
) -> str:
    rows = _collect_cluster_summaries(gene_dir)
    if not rows:
        return "No cluster summaries found in TSV files."
    if top_n and top_n > 0:
        rows = rows[:top_n]
    lines: List[str] = ["Most common structure + transcript by cluster ID:"]
    for row in rows:
        lines.append(
            f"{row['cluster']}: {row['structure']}; transcript_id={row['transcript_id']} "
            f"(structure_n={row['structure_n']}, tx_n={row['tx_n']})"
        )
    return "\n".join(lines)


def _add_picture_fit(slide, image_path: str, left, top, width, height) -> None:
    pic = slide.shapes.add_picture(image_path, left, top)
    scale = min(float(width) / float(pic.width), float(height) / float(pic.height))
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(left + (width - pic.width) / 2)
    pic.top = int(top + (height - pic.height) / 2)


def _find_existing_gene_image(image_dir: str, gene: str, preferred_format: str) -> Optional[str]:
    preferred = os.path.join(image_dir, f"{gene}.{preferred_format}")
    if os.path.exists(preferred) and os.path.getsize(preferred) > 0:
        return preferred
    for ext in ("png", "jpg"):
        if ext == preferred_format:
            continue
        alt = os.path.join(image_dir, f"{gene}.{ext}")
        if os.path.exists(alt) and os.path.getsize(alt) > 0:
            return alt
    return None


def build_pptx_from_isv(
    results_dir: str,
    output_dir: str,
    pptx_name: str,
    combined_glob: str = "combined__*.pdf",
    jpg_dpi: int = 220,
    image_format: str = "jpg",
    top_n_isoforms: int = 0,
    skip_existing_images: bool = False,
    summary_tsv_name: Optional[str] = None,
    event_table_tsv: Optional[str] = None,
    max_events_per_gene: int = 3,
    unsorted: bool = False,
    exclude_low_count_genes: int = 0,
) -> str:
    Presentation, Inches, Pt = _import_pptx()

    gene_inputs = find_gene_inputs(results_dir, combined_glob=combined_glob)
    if not gene_inputs:
        raise RuntimeError(
            f"No gene directories with '{combined_glob}' found under: {results_dir}"
        )

    os.makedirs(output_dir, exist_ok=True)
    image_format = str(image_format).lower()
    if image_format not in {"jpg", "png"}:
        raise ValueError("image_format must be 'jpg' or 'png'")
    image_dir = os.path.join(output_dir, f"isv_{image_format}")
    os.makedirs(image_dir, exist_ok=True)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    filtered_gene_inputs = gene_inputs
    if exclude_low_count_genes and exclude_low_count_genes > 0:
        filtered_gene_inputs = [
            item
            for item in gene_inputs
            if _passes_min_reads_any_condition(item.gene_dir, exclude_low_count_genes)
        ]
        if not filtered_gene_inputs:
            raise RuntimeError(
                f"All genes were excluded by --exclude_low_count_genes {exclude_low_count_genes}."
            )

    scored_items: List[Tuple[GeneInput, GeneScore]] = [
        (item, compute_gene_shift_score(item.gene_dir)) for item in filtered_gene_inputs
    ]
    if unsorted:
        scored_items.sort(key=lambda x: x[0].gene)
    else:
        scored_items.sort(key=lambda x: x[1].score, reverse=True)
    tsv_rows: List[Dict[str, Any]] = []
    event_rows_by_gene: Dict[str, List[Dict[str, Any]]] = {}
    if event_table_tsv:
        event_rows_by_gene, _ = _load_event_rows(event_table_tsv)

    for rank, (item, gene_score) in enumerate(scored_items, start=1):
        preferred_image_path = os.path.join(image_dir, f"{item.gene}.{image_format}")
        image_path = preferred_image_path
        existing_image_path = None
        if skip_existing_images:
            existing_image_path = _find_existing_gene_image(image_dir, item.gene, image_format)
        if existing_image_path:
            image_path = existing_image_path
        else:
            convert_pdf_to_image(item.combined_pdf, preferred_image_path, dpi=jpg_dpi, image_format=image_format)
            image_path = preferred_image_path
        cluster_rows = _collect_cluster_summaries(item.gene_dir)
        if top_n_isoforms and top_n_isoforms > 0:
            bottom_rows = cluster_rows[:top_n_isoforms]
        else:
            bottom_rows = cluster_rows
        if bottom_rows:
            lines = ["Most common structure + transcript by cluster ID:"]
            for row in bottom_rows:
                lines.append(
                    f"{row['cluster']}: {row['structure']}; transcript_id={row['transcript_id']} "
                    f"(structure_n={row['structure_n']}, tx_n={row['tx_n']})"
                )
            bottom_summary = "\n".join(lines)
        else:
            bottom_summary = "No cluster summaries found in TSV files."
        matched_event_rows = event_rows_by_gene.get(item.gene, [])
        if matched_event_rows:
            isoform_cluster_lookup = _build_isoform_cluster_lookup(item.gene_dir)
            event_summary = _render_event_summary(
                matched_event_rows,
                isoform_cluster_lookup,
                max_events_per_gene=max_events_per_gene,
            )
            if event_summary:
                bottom_summary = f"{bottom_summary}\n{event_summary}"

        slide = prs.slides.add_slide(blank_layout)
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Title near top, centered.
        title_box = slide.shapes.add_textbox(
            Inches(0.25), Inches(0.05), slide_w - Inches(0.5), Inches(0.35)
        )
        title_tf = title_box.text_frame
        title_tf.clear()
        p = title_tf.paragraphs[0]
        unique_exon_label = _short_unique_exon_label(item.gene_dir)
        if unique_exon_label:
            p.text = f"{item.gene} - {unique_exon_label}"
        else:
            p.text = item.gene
        title_text = p.text
        p.alignment = 1  # center
        p.font.size = Pt(20)

        score_box = slide.shapes.add_textbox(
            Inches(0.25), Inches(0.32), slide_w - Inches(0.5), Inches(0.22)
        )
        score_tf = score_box.text_frame
        score_tf.clear()
        sp = score_tf.paragraphs[0]
        sp.alignment = 1  # center
        sp.font.size = Pt(11)
        sp.text = (
            f"Shift score={gene_score.score:.3f} | "
            f"raw_ratio_shift={gene_score.raw_ratio_shift:.3f} | "
            f"min condition reads={int(gene_score.min_condition_reads)} | "
            f"read_weight={gene_score.read_weight:.3f} | "
            f"weight_factor={gene_score.weight_factor:.3f} | "
            f"major clusters={gene_score.major_clusters[0]},{gene_score.major_clusters[1]} | "
            f"cluster shifts={gene_score.cluster_shifts[0]:.3f},{gene_score.cluster_shifts[1]:.3f}"
        )

        # Main image region.
        img_left = Inches(0.2)
        img_top = Inches(0.58)
        img_w = slide_w - Inches(0.4)
        img_h = slide_h - Inches(1.68)
        _add_picture_fit(slide, image_path, img_left, img_top, img_w, img_h)

        # Bottom summary text.
        bottom_box = slide.shapes.add_textbox(
            Inches(0.2), slide_h - Inches(1.0), slide_w - Inches(0.4), Inches(0.9)
        )
        bottom_tf = bottom_box.text_frame
        bottom_tf.clear()
        p2 = bottom_tf.paragraphs[0]
        p2.text = bottom_summary
        p2.font.size = Pt(11)

        top1 = cluster_rows[0] if len(cluster_rows) > 0 else {}
        top2 = cluster_rows[1] if len(cluster_rows) > 1 else {}
        tsv_rows.append(
            {
                "slide_rank": rank,
                "gene": item.gene,
                "title_text": title_text,
                "unique_exon_label": unique_exon_label,
                "shift_score": round(gene_score.score, 6),
                "raw_ratio_shift": round(gene_score.raw_ratio_shift, 6),
                "min_condition_reads": int(gene_score.min_condition_reads),
                "read_weight": round(gene_score.read_weight, 6),
                "weight_factor": round(gene_score.weight_factor, 6),
                "major_cluster_1": gene_score.major_clusters[0],
                "major_cluster_2": gene_score.major_clusters[1],
                "cluster_shift_1": round(gene_score.cluster_shifts[0], 6),
                "cluster_shift_2": round(gene_score.cluster_shifts[1], 6),
                "top_cluster_1_label": top1.get("cluster", ""),
                "top_cluster_1_structure": top1.get("structure", ""),
                "top_cluster_1_transcript_id": top1.get("transcript_id", ""),
                "top_cluster_1_structure_n": top1.get("structure_n", ""),
                "top_cluster_1_tx_n": top1.get("tx_n", ""),
                "top_cluster_2_label": top2.get("cluster", ""),
                "top_cluster_2_structure": top2.get("structure", ""),
                "top_cluster_2_transcript_id": top2.get("transcript_id", ""),
                "top_cluster_2_structure_n": top2.get("structure_n", ""),
                "top_cluster_2_tx_n": top2.get("tx_n", ""),
                "matched_event_rows": len(matched_event_rows),
                "bottom_summary": bottom_summary,
                "combined_pdf": item.combined_pdf,
                "image_path": image_path,
            }
        )

    out_path = os.path.join(output_dir, pptx_name)
    if not out_path.lower().endswith(".pptx"):
        out_path += ".pptx"
    prs.save(out_path)
    tsv_base = summary_tsv_name
    if not tsv_base:
        ppt_base = os.path.splitext(os.path.basename(out_path))[0]
        tsv_base = f"{ppt_base}.tsv"
    tsv_path = os.path.join(output_dir, tsv_base)
    if not tsv_path.lower().endswith(".tsv"):
        tsv_path += ".tsv"
    fieldnames = list(tsv_rows[0].keys()) if tsv_rows else [
        "slide_rank", "gene", "title_text", "bottom_summary", "combined_pdf", "image_path"
    ]
    with open(tsv_path, "w", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames, delimiter="\t")
        writer.writeheader()
        for row in tsv_rows:
            writer.writerow(row)
    print(f"[summary_ISV] wrote TSV: {tsv_path}")
    return out_path


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Convert ISV combined PDFs to images and summarize them in a PowerPoint."
    )
    parser.add_argument("results_dir", help="Directory containing per-gene ISV output subdirectories.")
    parser.add_argument("output_dir", help="Directory where JPGs and PPTX will be written.")
    parser.add_argument("pptx_name", help="PowerPoint file name (with or without .pptx).")
    parser.add_argument(
        "--combined-glob",
        default="combined__*.pdf",
        help="Glob pattern for per-gene combined PDF (default: combined__*.pdf).",
    )
    parser.add_argument(
        "--jpg-dpi",
        type=int,
        default=220,
        help="DPI used for PDF->image conversion (default: 220).",
    )
    parser.add_argument(
        "--image-format",
        choices=["jpg", "png"],
        default="jpg",
        help="Intermediate image format embedded into PPTX (default: jpg).",
    )
    parser.add_argument(
        "--top-n-isoforms",
        type=int,
        default=0,
        help="Max number of cluster summaries at slide bottom (default: 0 = all clusters).",
    )
    parser.add_argument(
        "--skip-existing-images",
        action="store_true",
        help="Reuse existing per-gene rendered images in output_dir/isv_<format> when present.",
    )
    parser.add_argument(
        "--summary-tsv-name",
        default=None,
        help="Optional output TSV file name (default: <pptx_basename>.tsv).",
    )
    parser.add_argument(
        "--unsorted",
        action="store_true",
        help="Do not sort by shift score; sort slides alphabetically by gene.",
    )
    parser.add_argument(
        "--exclude_low_count_genes",
        type=int,
        default=0,
        metavar="MIN_READS",
        help=(
            "Exclude genes where no condition reaches at least MIN_READS total reads "
            "in *_isoform_ids.tsv files."
        ),
    )
    parser.add_argument(
        "--event-table-tsv",
        default=None,
        help=(
            "Optional significant results TSV with columns including "
            "Gene_Symbol, Feature, Coord_1/Coord_2, Isoform_1|Length, Isoform_2|Length, DiffMeans_*."
        ),
    )
    parser.add_argument(
        "--max-events-per-gene",
        type=int,
        default=3,
        help="Max number of matched event rows to print per gene at slide bottom (default: 3).",
    )
    return parser


def main(argv: Optional[List[str]] = None) -> int:
    parser = _build_parser()
    args = parser.parse_args(argv)
    out_path = build_pptx_from_isv(
        results_dir=args.results_dir,
        output_dir=args.output_dir,
        pptx_name=args.pptx_name,
        combined_glob=args.combined_glob,
        jpg_dpi=args.jpg_dpi,
        image_format=args.image_format,
        top_n_isoforms=args.top_n_isoforms,
        skip_existing_images=args.skip_existing_images,
        summary_tsv_name=args.summary_tsv_name,
        event_table_tsv=args.event_table_tsv,
        max_events_per_gene=args.max_events_per_gene,
        unsorted=args.unsorted,
        exclude_low_count_genes=args.exclude_low_count_genes,
    )
    print(f"[summary_ISV] wrote: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
